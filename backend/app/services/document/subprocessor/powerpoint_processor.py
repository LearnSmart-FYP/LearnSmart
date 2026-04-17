import logging
import base64
import shutil
import subprocess
import tempfile
from pathlib import Path
from uuid import UUID
from zipfile import ZipFile
from pptx import Presentation
import fitz  # PyMuPDF

from app.services.document.subprocessor.image_processor import image_processor
from app.services.document.subprocessor.table_processor import table_processor
from app.services.document.extraction_tree import ExtractionTree

logger = logging.getLogger(__name__)

def convert_to_pdf_for_ocr(file_path: str) -> str:

    file_path = Path(file_path)

    with tempfile.TemporaryDirectory() as temp_dir:

        profile_dir = tempfile.mkdtemp(prefix="lo_profile_")

        try:
            result = subprocess.run(
                [
                    'libreoffice',
                    '--headless',
                    f'-env:UserInstallation=file://{profile_dir}',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir,
                    str(file_path)
                ],
                capture_output = True,
                text = True,
                timeout = 180)
        finally:
            shutil.rmtree(profile_dir, ignore_errors = True)

        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice PDF conversion failed: {result.stderr}")

        converted_path = Path(temp_dir) / f"{file_path.stem}.pdf"

        if not converted_path.exists():
            raise RuntimeError(f"Converted PDF not found: {converted_path}")

        final_path = Path(tempfile.gettempdir()) / f"{file_path.stem}_ocr_{id(file_path)}.pdf"
        final_path.write_bytes(converted_path.read_bytes())

        return str(final_path)


def convert_ppt_to_pptx(ppt_path: str) -> str:

    ppt_path = Path(ppt_path)

    with tempfile.TemporaryDirectory() as temp_dir:
        result = subprocess.run(
            [
                'libreoffice',
                '--headless',
                '--convert-to', 'pptx',
                '--outdir', temp_dir,
                str(ppt_path)
            ],
            capture_output = True,
            text = True,
            timeout = 120)

        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        converted_path = Path(temp_dir) / f"{ppt_path.stem}.pptx"

        if not converted_path.exists():
            raise RuntimeError(f"Converted file not found: {converted_path}")

        final_path = ppt_path.with_suffix('.pptx')
        final_path.write_bytes(converted_path.read_bytes())

        logger.info(f"Converted {ppt_path.name} to {final_path.name}")
        return str(final_path)

class PowerPointProcessor:

    def __init__(self):

        # Use global helper processors
        self.image_processor = image_processor
        self.table_processor = table_processor

    def _get_slide_ocr_texts(self, file_path: str) -> dict[int, str]:

        ocr_texts = {}
        pdf_path = None

        try:

            logger.info(f"Converting to PDF for full-slide OCR")
            pdf_path = convert_to_pdf_for_ocr(file_path)
            doc = fitz.open(pdf_path)

            for i, page in enumerate(doc):
                slide_number = i + 1
                try:
                    mat = fitz.Matrix(2, 2)
                    pix = page.get_pixmap(matrix = mat)
                    page_image_bytes = pix.tobytes("png")

                    ocr_result = self.image_processor.process(page_image_bytes)
                    if ocr_result.success and ocr_result.output_text:
                        ocr_texts[slide_number] = ocr_result.output_text
                except Exception as e:
                    logger.debug(f"Full-slide OCR failed for slide {slide_number}: {e}")

            doc.close()
            logger.info(f"Full-slide OCR completed for {len(ocr_texts)} slides")

        except Exception as e:
            logger.warning(f"PDF conversion for OCR failed: {e}")

        finally:
            if pdf_path and Path(pdf_path).exists():
                try:
                    Path(pdf_path).unlink()
                except Exception:
                    pass

        return ocr_texts

    def _process_shape(
        self,
        shape,
        tree: ExtractionTree,
        parent_node_id: UUID,
        file_path: str,
        slide_number: int,
        doc_base_name: str,
        text_parts: list[str],
        title: str,
        counters: dict) -> None:

        # Handle text shapes (AUTO_SHAPE, CALLOUT, COMMENT, PLACEHOLDER, TEXT_BOX)
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text and text != title:
                text_parts.append(text)

        # Extract tables
        if shape.shape_type == 19:
            try:
                table = shape.table
                table_data = []

                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)

                if table_data:
                    table_result = self.table_processor.process(table_data = table_data)

                    if table_result.success:
                        tree.add_child(
                            parent_id = parent_node_id,
                            source_path = f"{file_path}:slide{slide_number}:table{counters['tables']}",
                            document_type = "table",
                            extraction_type = "table",
                            metadata = {
                                "output_text": table_result.output_text,
                                "suggested_filename": f"{doc_base_name}_slide{slide_number}_table{counters['tables']}.csv"})
                        counters['tables'] += 1

            except Exception as e:
                logger.debug(f"Table extraction failed on slide {slide_number}: {e}")

        # Extract images
        if shape.shape_type == 13:
            try:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext if hasattr(image, 'ext') else "png"

                output_text = ""
                try:
                    ocr_result = self.image_processor.process(image_bytes)
                    output_text = ocr_result.output_text if ocr_result.success else ""
                except Exception:
                    output_text = ""

                tree.add_child(
                    parent_id = parent_node_id,
                    source_path = f"{file_path}:slide{slide_number}:image{counters['images']}",
                    document_type = "image",
                    extraction_type = "image",
                    metadata = {
                        "data": base64.b64encode(image_bytes).decode('utf-8'),
                        "format": image_ext,
                        "output_text": output_text,
                        "extraction_location": f"slide {slide_number}",
                        "suggested_filename": f"{doc_base_name}_slide{slide_number}_img{counters['images']}.{image_ext}"})
                counters['images'] += 1

            except Exception as e:
                logger.debug(f"Image extraction failed on slide {slide_number}: {e}")

        # Extract charts (not supported for now)
        if shape.shape_type == 6:
            try:
                tree.add_child(
                    parent_id = parent_node_id,
                    source_path = f"{file_path}:slide{slide_number}:chart{counters['charts']}",
                    document_type = "chart",
                    extraction_type = "chart",
                    metadata = {
                        "extraction_location": f"slide {slide_number}",
                        "note": "Chart detected but image export not supported"})
                counters['charts'] += 1

            except Exception as e:
                logger.debug(f"Chart extraction failed on slide {slide_number}: {e}")

        # Extract group shapes
        if shape.shape_type == 24:
            try:
                if hasattr(shape, 'shapes'):
                    for child_shape in shape.shapes:
                        self._process_shape(
                            shape = child_shape,
                            tree = tree,
                            parent_node_id = parent_node_id,
                            file_path = file_path,
                            slide_number = slide_number,
                            doc_base_name = doc_base_name,
                            text_parts = text_parts,
                            title = title,
                            counters = counters)
            except Exception as e:
                logger.debug(f"Group extraction failed on slide {slide_number}: {e}")

    def process_sync(
        self,
        file_path: str,
        tree: ExtractionTree | None = None,
        parent_node_id: UUID | None = None) -> ExtractionTree:

        total_embedded = 0
        total_media = 0
        counters = {'tables': 0, 'images': 0, 'charts': 0}
        converted_file = None

        try:

            # Convert .ppt to .pptx if necessary

            file_path_obj = Path(file_path)
            if file_path_obj.suffix.lower() == '.ppt':
                logger.info(f"Converting legacy .ppt file: {file_path}")
                converted_file = convert_ppt_to_pptx(file_path)
                file_path = converted_file

            prs = Presentation(file_path)
            doc_base_name = Path(file_path).stem

            slide_ocr_texts = self._get_slide_ocr_texts(file_path)

            # Create or use existing tree

            if tree is None:
                tree = ExtractionTree()
                root_id = tree.add_root(
                    source_path = file_path,
                    document_type = "powerpoint",
                    metadata = {})
            else:
                root_id = tree.add_child(
                    parent_id = parent_node_id,
                    source_path = file_path,
                    document_type = "powerpoint",
                    extraction_type = "embedded_document",
                    metadata = {})

            # Process each slide

            for slide_idx, slide in enumerate(prs.slides):

                slide_number = slide_idx + 1
                title = ""
                text_parts = []

                # Add slide node
                slide_node_id = tree.add_child(
                    parent_id = root_id,
                    source_path = f"{file_path}:slide{slide_number}",
                    document_type = "slide",
                    extraction_type = "slide",
                    metadata = {"page_number": slide_number})

                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()

                # Process all shapes (including recursive group shapes)
                for shape in slide.shapes:
                    self._process_shape(
                        shape = shape,
                        tree = tree,
                        parent_node_id = slide_node_id,
                        file_path = file_path,
                        slide_number = slide_number,
                        doc_base_name = doc_base_name,
                        text_parts = text_parts,
                        title = title,
                        counters = counters)

                # Combine native text + full-slide OCR
                native_text = "\n".join(text_parts)
                ocr_text = slide_ocr_texts.get(slide_number, "")

                final_text_parts = []
                if native_text.strip():
                    final_text_parts.append(native_text.strip())
                if ocr_text.strip():
                    final_text_parts.append(f"\n[Full Slide OCR]:\n{ocr_text.strip()}")

                slide_text = "\n".join(final_text_parts)

                slide_node = tree.get_node(slide_node_id)
                if slide_node:
                    slide_node.metadata.update({
                        "output_text": slide_text})

            # Extract embedded files and media from presentation
            try:

                with ZipFile(file_path, 'r') as zf:

                    for name in zf.namelist():

                        if name.startswith('ppt/embeddings/') and not name.endswith('/'):
                            try:
                                emb_bytes = zf.read(name)
                                emb_filename = Path(name).name
                                safe_emb_name = f"{doc_base_name}_{emb_filename}"

                                tree.add_child(
                                    parent_id = root_id,
                                    source_path = f"{file_path}:embedded:{emb_filename}",
                                    document_type = "embedded",
                                    extraction_type = "embedded_document",
                                    metadata = {
                                        "data": base64.b64encode(emb_bytes).decode('utf-8'),
                                        "original_filename": emb_filename,
                                        "file_size": len(emb_bytes),
                                        "suggested_filename": safe_emb_name})

                                total_embedded += 1
                                logger.debug(f"Extracted embedded file: {safe_emb_name}")

                            except Exception as e:
                                logger.debug(f"Embedded file extraction failed for {name}: {e}")

                        if name.startswith('ppt/media/') and not name.endswith('/'):
                            try:
                                media_filename = Path(name).name
                                ext = Path(media_filename).suffix.lower()

                                if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf', '.emf']:
                                    continue

                                media_bytes = zf.read(name)
                                safe_media_name = f"{doc_base_name}_{media_filename}"
                                tree.add_child(
                                    parent_id = root_id,
                                    source_path = f"{file_path}:media:{media_filename}",
                                    document_type = "media",
                                    extraction_type = "media",
                                    metadata = {
                                        "data": base64.b64encode(media_bytes).decode('utf-8'),
                                        "original_filename": media_filename,
                                        "file_size": len(media_bytes),
                                        "suggested_filename": safe_media_name})

                                total_media += 1
                                logger.debug(f"Extracted media file: {safe_media_name}")

                            except Exception as e:
                                logger.debug(f"Media file extraction failed for {name}: {e}")

            except Exception as e:
                logger.debug(f"Embedded/media files scan failed: {e}")

            logger.info(f"PowerPoint processed: {len(prs.slides)} slides, "
                       f"Tables: {counters['tables']}, Images: {counters['images']}, "
                       f"Media: {total_media}, Embedded: {total_embedded}, "
                       f"Slides with OCR: {len(slide_ocr_texts)}")

            return tree

        except Exception as e:
            logger.error(f"PowerPoint processing failed: {e}")
            raise

        finally:

            if converted_file and Path(converted_file).exists():
                try:
                    Path(converted_file).unlink()
                except Exception:
                    pass

# Global instance
powerpoint_processor = PowerPointProcessor()
